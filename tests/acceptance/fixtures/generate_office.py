"""Generate PPTX and XLSX acceptance fixtures."""
from pathlib import Path

_HERE = Path(__file__).resolve().parent
_DEST = _HERE / "office"
_DEST.mkdir(parents=True, exist_ok=True)

# ── PPTX ──
from pptx import Presentation

prs = Presentation()
prs.slide_width = 9144000  # 10 inches
prs.slide_height = 6858000  # 7.5 inches

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "OKS PPTX Acceptance Test"
slide.shapes.placeholders[1].text = "v0.4.0 — Multi-slide with tables and images"

# Slide 2: Content with table
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Capability Matrix"
rows, cols = 4, 3
table = slide.shapes.add_table(rows, cols, 1000000, 1500000, 6000000, 3000000).table
headers = ["Provider", "Capability", "Maturity"]
for c, h in enumerate(headers):
    table.cell(0, c).text = h
data = [
    ["pdf-lite", "document.text.extract", "validated"],
    ["rapidocr", "image.ocr", "validated"],
    ["firecrawl", "web.extract", "validated"],
]
for r, row_data in enumerate(data, 1):
    for c, val in enumerate(row_data):
        table.cell(r, c).text = val

# Slide 3: Bullet points
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Known Limitations"
body = slide.shapes.placeholders[1].text_frame
body.text = "PPTX acceptance covers:"
for item in [
    "Slide text extraction with page locator",
    "Table structure preservation",
    "Multi-slide documents",
]:
    p = body.add_paragraph()
    p.text = f"• {item}"
    p.level = 1

# Slide 4: Chart placeholder (no actual chart data — honesty test)
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Chart Placeholder"
body = slide.shapes.placeholders[1].text_frame
body.text = "This slide references a chart image."
p = body.add_paragraph()
p.text = "Chart interpretation requires image.observe or chart.interpret capability."
p = body.add_paragraph()
p.text = "(Chart image not embedded — this is a partial test case.)"

pptx_path = _DEST / "acceptance.pptx"
prs.save(str(pptx_path))
print(f"PPTX: {pptx_path} ({pptx_path.stat().st_size} bytes, {len(prs.slides)} slides)")

# ── XLSX ──
from openpyxl import Workbook

wb = Workbook()

# Sheet 1: Basic data
ws1 = wb.active
ws1.title = "Metrics"
ws1.append(["Metric", "Value", "Unit"])
ws1.append(["Evidence count", 46, "records"])
ws1.append(["Text chars", 696, "chars"])
ws1.append(["Pages", 3, "pages"])
ws1.append(["Latency", 6.2, "seconds"])
ws1.merge_cells("A6:C6")
ws1["A6"] = "Merged cell: summary row"

# Sheet 2: With formula
ws2 = wb.create_sheet("Formulas")
ws2.append(["Item", "Qty", "Price", "Total"])
ws2.append(["Widget", 10, 2.5, None])
ws2.append(["Gadget", 5, 8.0, None])
ws2["D2"] = "=B2*C2"
ws2["D3"] = "=B3*C3"
ws2["D4"] = "=SUM(D2:D3)"
ws2["A5"] = "Formula cells above — formula_preserved: true, formula_evaluated: false"

# Sheet 3: Empty
wb.create_sheet("Empty")

xlsx_path = _DEST / "acceptance.xlsx"
wb.save(str(xlsx_path))
print(f"XLSX: {xlsx_path} ({xlsx_path.stat().st_size} bytes, {len(wb.sheetnames)} sheets: {wb.sheetnames})")
